"""Build the Tools & Libraries presentation.

Creates docs/Tools_and_Libraries_Presentation.pptx with:
  - a title slide
  - a project overview slide
  - an architecture / how-it-works slide
  - a tech-stack grid
  - one image slide per tool / library (short description)
  - a summary slide

Run docs/make_tools_images.py first to (re)generate the icon images.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DOCS = Path(__file__).parent
IMG = DOCS / "tools_images"

DARK = RGBColor(0x1E, 0x29, 0x3B)
ACCENT = RGBColor(0x38, 0x7A, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x4A, 0x55, 0x68)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def bar(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


# ---------------- Title slide ----------------
s = slide()
bg(s, DARK)
textbox(s, Inches(0.6), Inches(2.35), Inches(12.1), Inches(1.2),
        [("Student Support Service Chatbot", 40, True, WHITE)], PP_ALIGN.CENTER)
bar(s, Inches(4.67), Inches(3.55), Inches(4), Inches(0.05), ACCENT)
textbox(s, Inches(1), Inches(3.75), Inches(11.3), Inches(1.2),
        [("Tools & Libraries Used", 26, False, RGBColor(0x9E, 0xC5, 0xF5))], PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.6),
        [("An AI-powered student help desk \u2022 Flask + scikit-learn + Groq LLM", 15, False, RGBColor(0xB5, 0xBD, 0xC9))],
        PP_ALIGN.CENTER)

# ---------------- Project overview ----------------
s = slide()
bg(s, WHITE)
bar(s, 0, 0, SW, Inches(1.1), DARK)
textbox(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
        [("Project Overview", 32, True, WHITE)])
overview = [
    ("What it is", 20, True, ACCENT),
    ("A web-based student support assistant that answers questions about college services, "
     "academics and internships through a simple chat interface.", 16, False, GREY),
    ("", 8, False, GREY),
    ("How it answers  \u2013  three layers", 20, True, ACCENT),
    ("\u2022  ML intent model (scikit-learn TF-IDF + Logistic Regression) for common queries", 16, False, GREY),
    ("\u2022  Rule-based regex patterns as a fast fallback", 16, False, GREY),
    ("\u2022  Groq LLM (llama-3.1) for Academic & Internship modes, grounded on local JSON data", 16, False, GREY),
    ("", 8, False, GREY),
    ("Chat modes", 20, True, ACCENT),
    ("General  \u2022  Academic (university data)  \u2022  Internship (live-style listings)", 16, False, GREY),
    ("", 8, False, GREY),
    ("Extras", 20, True, ACCENT),
    ("Session login, browser-stored API key, and a unit-test suite.", 16, False, GREY),
]
textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.9), overview)

# ---------------- Architecture / flow ----------------
s = slide()
bg(s, WHITE)
bar(s, 0, 0, SW, Inches(1.1), DARK)
textbox(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
        [("How It Works", 32, True, WHITE)])
from pptx.enum.shapes import MSO_SHAPE


def flowbox(x, title, subtitle, color):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.7), Inches(2.3), Inches(1.7))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    tf = sp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = subtitle; r2.font.size = Pt(11); r2.font.color.rgb = WHITE


def arrow(x):
    sp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.25), Inches(0.55), Inches(0.6))
    sp.fill.solid(); sp.fill.fore_color.rgb = ACCENT; sp.line.fill.background()


xs = [Inches(0.55), Inches(3.55), Inches(6.55), Inches(9.55)]
flowbox(xs[0], "Browser UI", "HTML / CSS / JS", RGBColor(0x38, 0x7A, 0xE0))
arrow(Inches(2.95))
flowbox(xs[1], "Flask API", "/api/chat, login", RGBColor(0x2E, 0x7D, 0x32))
arrow(Inches(5.95))
flowbox(xs[2], "Chatbot Engine", "ML \u2192 rules \u2192 LLM", RGBColor(0xF4, 0x5B, 0x42))
arrow(Inches(8.95))
flowbox(xs[3], "Response", "text reply", RGBColor(0x5E, 0x35, 0xB1))
textbox(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.8),
        [("Groq LLM API and JSON knowledge files (university_data.json, internship_template.json) "
          "power the Academic and Internship modes; joblib loads the pre-trained scikit-learn pipeline.",
          15, False, GREY)])

# ---------------- Tech stack grid ----------------
s = slide()
bg(s, WHITE)
bar(s, 0, 0, SW, Inches(1.1), DARK)
textbox(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
        [("Tech Stack at a Glance", 32, True, WHITE)])
groups = [
    ("Backend", ["Python", "Flask", "requests"], RGBColor(0x2E, 0x7D, 0x32)),
    ("Machine Learning", ["scikit-learn", "joblib"], RGBColor(0xF4, 0x8C, 0x1D)),
    ("AI / LLM", ["Groq API (llama-3.1)"], RGBColor(0xF4, 0x5B, 0x42)),
    ("Frontend", ["HTML5", "CSS3", "JavaScript"], RGBColor(0x38, 0x7A, 0xE0)),
    ("Data", ["JSON files"], RGBColor(0x5E, 0x35, 0xB1)),
    ("Testing", ["unittest"], RGBColor(0x00, 0x83, 0x8F)),
]
cols = 3
cw, ch = Inches(3.9), Inches(2.4)
gx, gy = Inches(0.7), Inches(1.5)
for i, (title, items, color) in enumerate(groups):
    r, c = divmod(i, cols)
    x = gx + c * (cw + Inches(0.25))
    y = gy + r * (ch + Inches(0.25))
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = color; card.line.width = Pt(2)
    tf = card.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]; rr = p.add_run(); rr.text = title; rr.font.bold = True; rr.font.size = Pt(20); rr.font.color.rgb = color
    for it in items:
        pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = "\u2022  " + it
        rr.font.size = Pt(15); rr.font.color.rgb = GREY

# ---------------- Per-tool slides ----------------
TOOLS = [
    ("python", "Python", "Core programming language",
     ["Runs the backend server and the ML training script",
      "Standard library: re, json, random, datetime for logic & data handling"]),
    ("flask", "Flask", "Lightweight web framework",
     ["Serves the static frontend and JSON REST API",
      "Routes: /api/login, /api/logout, /api/session, /api/chat",
      "Uses server-side sessions for login state"]),
    ("scikit_learn", "scikit-learn", "Machine learning toolkit",
     ["TfidfVectorizer + LogisticRegression in a Pipeline",
      "Classifies a message into an intent (greeting, hours, registration...)",
      "Trained by train_model.py"]),
    ("joblib", "joblib", "Model persistence",
     ["Saves & loads the trained pipeline and responses",
      "Files: model/pipeline.joblib, model/responses.joblib",
      "Lets the app skip re-training at startup"]),
    ("requests", "requests", "HTTP client library",
     ["Calls the Groq chat-completions REST endpoint",
      "Sends the API key as a Bearer token with a 30s timeout",
      "Handles HTTP errors (rate limit, auth) gracefully"]),
    ("groq", "Groq LLM API", "Large language model service",
     ["Powers Academic & Internship modes (model: llama-3.1-8b-instant)",
      "Prompts are grounded on local university & internship JSON data",
      "Key supplied via env var or saved in the browser"]),
    ("html5", "HTML5", "Page structure",
     ["index.html builds the login panel and chat UI",
      "Semantic markup with accessible labels & aria attributes"]),
    ("css3", "CSS3", "Styling & layout",
     ["styles.css themes the login card, chat bubbles and mode buttons",
      "Responsive, modern look with no CSS framework"]),
    ("javascript", "JavaScript (Vanilla)", "Frontend logic",
     ["app.js handles login, chat, typing indicator and mode switching",
      "Uses the Fetch API to talk to the Flask backend",
      "Stores the Groq API key in localStorage"]),
    ("json", "JSON Data", "Knowledge files",
     ["university_data.json grounds Academic answers",
      "internship_template.json feeds the live-style internship feed",
      "Loaded and injected into LLM prompts"]),
    ("unittest", "unittest", "Automated testing",
     ["tests/ verify chatbot replies and login-protected chat",
      "Mocks Groq calls and checks the API key is never leaked"]),
]

for fname, name, tagline, bullets in TOOLS:
    s = slide()
    bg(s, WHITE)
    bar(s, 0, 0, Inches(5.6), SH, DARK)  # left panel
    img_path = IMG / f"{fname}.png"
    if img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(0.55), Inches(2.35), width=Inches(4.5))
    # right content
    textbox(s, Inches(6.1), Inches(1.3), Inches(6.7), Inches(1.2),
            [(name, 38, True, DARK)])
    textbox(s, Inches(6.1), Inches(2.3), Inches(6.7), Inches(0.8),
            [(tagline, 20, True, ACCENT)])
    lines = [("\u2022  " + b, 17, False, GREY) for b in bullets]
    textbox(s, Inches(6.1), Inches(3.2), Inches(6.7), Inches(3.5), lines)

# ---------------- Summary ----------------
s = slide()
bg(s, DARK)
textbox(s, Inches(1), Inches(0.9), Inches(11.3), Inches(1.0),
        [("Summary", 40, True, WHITE)], PP_ALIGN.CENTER)
summary = [
    ("A compact full-stack AI app built from focused, well-known tools:", 20, False, RGBColor(0xB5, 0xBD, 0xC9)),
    ("", 10, False, WHITE),
    ("\u2022  Python + Flask  \u2013  backend & REST API", 19, False, WHITE),
    ("\u2022  scikit-learn + joblib  \u2013  ML intent classification", 19, False, WHITE),
    ("\u2022  requests + Groq LLM  \u2013  smart, grounded answers", 19, False, WHITE),
    ("\u2022  HTML / CSS / JavaScript  \u2013  the chat interface", 19, False, WHITE),
    ("\u2022  JSON data + unittest  \u2013  knowledge & quality", 19, False, WHITE),
    ("", 10, False, WHITE),
    ("Simple to run:  python app.py  \u2192  http://127.0.0.1:5000/", 18, True, RGBColor(0x9E, 0xC5, 0xF5)),
]
textbox(s, Inches(2.2), Inches(2.2), Inches(9), Inches(4.8), summary)

out = DOCS / "Tools_and_Libraries_Presentation.pptx"
prs.save(out)
print("Saved", out)
